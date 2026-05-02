"""
pdf_to_pptx.py — Convert a PDF (one slide per page) into a PPTX file where
each page is placed as a full-slide JPEG image.

Usage (single file):
    uv run python pdf_to_pptx.py presentation.pdf
    uv run python pdf_to_pptx.py presentation.pdf --output out.pptx

Usage (batch directory):
    uv run python pdf_to_pptx.py ./slides_dir --batch true
    uv run python pdf_to_pptx.py ./slides_dir --batch true --output ./out_dir

Usage (with speaker notes from Keynote):
    uv run python pdf_to_pptx.py presentation.pdf --with-notes true
    uv run python pdf_to_pptx.py presentation.pdf --with-notes true --keynote "~/Decks/My Talk.key"
    uv run python pdf_to_pptx.py ./slides_dir --batch true --with-notes true

When --with-notes is used the script exports a temporary PPTX from the
matching open Keynote document, reads click-build animation counts to figure
out which Keynote slide maps to which PDF page(s), and copies the speaker
notes into the output PPTX.  Use --keynote to specify the Keynote file
explicitly instead of being prompted.
"""

from __future__ import annotations

import io
import tempfile
from enum import StrEnum
from pathlib import Path

import fitz  # PyMuPDF
from jsonargparse import auto_cli
from pptx import Presentation
from pptx.util import Emu
from pydantic import BaseModel, Field
from rich.console import Console
from rich.progress import BarColumn, Progress, SpinnerColumn, TaskID, TaskProgressColumn, TextColumn
from rich.table import Table

console = Console()

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

POINTS_TO_EMU: int = 12700  # 1 pt = 12 700 EMU


# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------


class ConversionStatus(StrEnum):
    SUCCESS = "success"
    FAILURE = "failure"


class NotesMismatchWarning(StrEnum):
    NONE = "none"
    COUNT_MISMATCH = "count_mismatch"


class SlideNotesMapping(BaseModel):
    keynote_slide_index: int = Field(description="0-based index in the temp PPTX (includes hidden slides)")
    pdf_page_start: int = Field(description="0-based index of first PDF page for this slide")
    pdf_page_end: int = Field(description="0-based index of last PDF page for this slide (inclusive)")
    notes_paragraphs: list[str] = Field(
        description="List of paragraph text strings from the notes (one string per paragraph, preserving empty strings for blank lines)"
    )


class NotesExtractionResult(BaseModel):
    mappings: list[SlideNotesMapping] = Field(description="Per-slide notes mappings")
    keynote_slide_count: int = Field(description="Total number of slides in the Keynote document (including hidden)")
    expected_pdf_pages: int = Field(description="Expected number of PDF pages based on animation build counts")
    actual_pdf_pages: int = Field(description="Actual number of pages in the PDF")
    warning: NotesMismatchWarning = Field(description="Warning about notes mapping mismatch, if any")
    warning_message: str | None = Field(
        default=None, description="Human-readable explanation of the mismatch, if any"
    )


class ConversionResult(BaseModel):
    input_path: Path = Field(description="Path to the source PDF file")
    output_path: Path | None = Field(description="Path to the produced PPTX file, or None on failure")
    page_count: int = Field(description="Number of pages converted (0 on failure)")
    status: ConversionStatus = Field(description="Whether the conversion succeeded or failed")
    error: str | None = Field(default=None, description="Error message if the conversion failed")
    notes_warning: NotesMismatchWarning | None = Field(
        default=None, description="Warning about notes mapping mismatch, if notes were extracted"
    )


# ---------------------------------------------------------------------------
# Core conversion logic
# ---------------------------------------------------------------------------


def convert_pdf_to_pptx(
    pdf_path: Path,
    output_path: Path,
    dpi: int,
    jpeg_quality: int,
    progress: Progress | None = None,
    progress_task: TaskID | None = None,
    notes_mapping: NotesExtractionResult | None = None,
) -> ConversionResult:
    """Convert a single PDF file to a PPTX file.

    Each PDF page is rendered to a JPEG at *dpi* resolution and placed as a
    full-bleed image on its own blank slide.  Slide dimensions are taken
    directly from the PDF page dimensions so aspect ratio is always preserved.

    If *notes_mapping* is provided, speaker notes from the matched Keynote
    slide are embedded into each output slide.
    """
    try:
        doc = fitz.open(str(pdf_path))
    except Exception as exc:  # noqa: BLE001
        return ConversionResult(
            input_path=pdf_path,
            output_path=None,
            page_count=0,
            status=ConversionStatus.FAILURE,
            error=f"Failed to open PDF: {exc}",
        )

    page_count = len(doc)
    if page_count == 0:
        return ConversionResult(
            input_path=pdf_path,
            output_path=None,
            page_count=0,
            status=ConversionStatus.FAILURE,
            error="PDF contains no pages",
        )

    try:
        prs = Presentation()
        scale = dpi / 72.0
        matrix = fitz.Matrix(scale, scale)

        for page_index in range(page_count):
            page = doc[page_index]

            # --- slide dimensions from PDF page (points → EMU) ---
            rect = page.rect  # in points
            slide_width = Emu(int(rect.width * POINTS_TO_EMU))
            slide_height = Emu(int(rect.height * POINTS_TO_EMU))

            # Set presentation dimensions from the first page.
            # (All pages will use the same declared size; if pages genuinely
            # differ in size this is the best we can do in PPTX.)
            if page_index == 0:
                prs.slide_width = slide_width
                prs.slide_height = slide_height

            # --- render page to JPEG bytes ---
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            jpeg_bytes = pixmap.tobytes("jpeg", jpg_quality=jpeg_quality)

            # --- add blank slide and fill with the image ---
            blank_layout = prs.slide_layouts[6]  # completely blank layout
            slide = prs.slides.add_slide(blank_layout)

            image_stream = io.BytesIO(jpeg_bytes)
            slide.shapes.add_picture(
                image_stream,
                left=Emu(0),
                top=Emu(0),
                width=slide_width,
                height=slide_height,
            )

            # --- copy speaker notes if a mapping is provided ---
            if notes_mapping is not None:
                matched: SlideNotesMapping | None = None
                for m in notes_mapping.mappings:
                    if m.pdf_page_start <= page_index <= m.pdf_page_end:
                        matched = m
                        break

                if matched is not None and matched.notes_paragraphs:
                    try:
                        tf = slide.notes_slide.notes_text_frame
                        tf.text = ""  # clear any existing content
                        paragraphs = matched.notes_paragraphs
                        tf.paragraphs[0].text = paragraphs[0]
                        for p in paragraphs[1:]:
                            tf.add_paragraph().text = p
                    except Exception:  # noqa: BLE001
                        pass  # notes are non-critical; skip silently

            if progress is not None and progress_task is not None:
                progress.advance(progress_task)

        doc.close()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        notes_warning: NotesMismatchWarning | None = None
        if notes_mapping is not None and notes_mapping.warning != NotesMismatchWarning.NONE:
            notes_warning = notes_mapping.warning

        return ConversionResult(
            input_path=pdf_path,
            output_path=output_path,
            page_count=page_count,
            status=ConversionStatus.SUCCESS,
            notes_warning=notes_warning,
        )

    except Exception as exc:  # noqa: BLE001
        doc.close()
        return ConversionResult(
            input_path=pdf_path,
            output_path=None,
            page_count=0,
            status=ConversionStatus.FAILURE,
            error=str(exc),
        )


# ---------------------------------------------------------------------------
# Output path resolution
# ---------------------------------------------------------------------------


def resolve_output_path(pdf_path: Path, output: Path | None, batch: bool) -> Path:
    """Determine where to write the PPTX for *pdf_path*."""
    if output is None:
        return pdf_path.with_suffix(".pptx")
    if batch:
        # output is treated as a directory
        return output / pdf_path.with_suffix(".pptx").name
    return output


# ---------------------------------------------------------------------------
# Validation helpers
# ---------------------------------------------------------------------------


def validate_inputs(input: Path, batch: bool) -> list[Path]:  # noqa: A002
    """Return the list of PDF paths to process, or raise ValueError."""
    if batch:
        if not input.is_dir():
            raise ValueError(f"--batch requires a directory, but got: {input}")
        pdfs = sorted(input.glob("*.pdf"))
        if not pdfs:
            raise ValueError(f"No .pdf files found in directory: {input}")
        return pdfs
    else:
        if not input.exists():
            raise ValueError(f"Input file not found: {input}")
        if input.suffix.lower() != ".pdf":
            raise ValueError(f"Input file is not a PDF: {input}")
        return [input]


# ---------------------------------------------------------------------------
# Summary rendering
# ---------------------------------------------------------------------------


def print_summary(results: list[ConversionResult]) -> None:
    """Render a Rich summary table of all conversion results."""
    successes = [r for r in results if r.status == ConversionStatus.SUCCESS]
    failures = [r for r in results if r.status == ConversionStatus.FAILURE]

    # Determine whether any result has notes info to show
    any_notes = any(r.notes_warning is not None for r in results)

    table = Table(title="Conversion Summary", show_header=True, header_style="bold cyan", expand=True)
    table.add_column("File", style="dim", ratio=3, no_wrap=True)
    table.add_column("Pages", justify="right", min_width=5)
    table.add_column("Output", ratio=5)
    table.add_column("Status", min_width=10)
    if any_notes:
        table.add_column("Notes", min_width=12)

    for r in results:
        if r.status == ConversionStatus.SUCCESS:
            status_cell = "[green]✓ success[/green]"
            output_cell = str(r.output_path)
            pages_cell = str(r.page_count)
        else:
            status_cell = "[red]✗ failed[/red]"
            output_cell = f"[red]{r.error}[/red]"
            pages_cell = "—"

        if any_notes:
            match r.notes_warning:
                case NotesMismatchWarning.COUNT_MISMATCH:
                    notes_cell = "[yellow]⚠ mismatch[/yellow]"
                case NotesMismatchWarning.NONE | None:
                    # NONE means notes were extracted cleanly; None means notes weren't requested
                    notes_cell = "[green]✓ embedded[/green]" if r.notes_warning == NotesMismatchWarning.NONE else "—"
            table.add_row(r.input_path.name, pages_cell, output_cell, status_cell, notes_cell)
        else:
            table.add_row(r.input_path.name, pages_cell, output_cell, status_cell)

    console.print()
    console.print(table)
    console.print(
        f"[bold]Done:[/bold] [green]{len(successes)} succeeded[/green]"
        + (f", [red]{len(failures)} failed[/red]" if failures else ""),
    )


# ---------------------------------------------------------------------------
# Notes extraction helper
# ---------------------------------------------------------------------------


def _extract_notes_for_pdf(
    pdf_path: Path,
    keynote_path: Path,
) -> NotesExtractionResult | None:
    """Export a temp PPTX from *keynote_path*, build the notes mapping for
    *pdf_path*, and return the result.  Returns None on any error (prints a
    warning to the console)."""
    # Import here to avoid circular-import issues and keep top-level clean
    from keynote import export_keynote_to_pptx  # type: ignore[import]
    from notes import build_notes_mapping  # type: ignore[import]

    try:
        doc = fitz.open(str(pdf_path))
        actual_page_count = len(doc)
        doc.close()
    except Exception as exc:  # noqa: BLE001
        console.print(f"[yellow]⚠ Could not open PDF to count pages for notes mapping: {exc}[/yellow]")
        return None

    tmp_pptx: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            tmp_pptx = Path(f.name)

        console.print(f"[dim]Exporting Keynote → temp PPTX: {keynote_path.name}[/dim]")
        export_keynote_to_pptx(keynote_path, tmp_pptx)

        mapping = build_notes_mapping(tmp_pptx, actual_page_count)

        if mapping.warning != NotesMismatchWarning.NONE:
            console.print(f"[yellow]⚠ Notes mapping warning: {mapping.warning_message}[/yellow]")

        return mapping

    except Exception as exc:  # noqa: BLE001
        console.print(f"[yellow]⚠ Notes extraction failed, continuing without notes: {exc}[/yellow]")
        return None
    finally:
        if tmp_pptx is not None and tmp_pptx.exists():
            try:
                tmp_pptx.unlink()
            except Exception:  # noqa: BLE001
                pass


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------


def main(
    input: Path,  # noqa: A002
    output: Path | None = None,
    dpi: int = 150,
    jpeg_quality: int = 85,
    batch: bool = False,
    with_notes: bool = False,
    keynote: Path | None = None,
) -> None:
    """Convert a PDF (or a directory of PDFs) into PPTX files.

    Each PDF page is rendered as a full-slide JPEG image at the specified DPI.
    Optionally, speaker notes can be pulled from an open Keynote document and
    embedded into the output PPTX.

    Args:
        input: Path to a PDF file, or a directory when --batch is True.
        output: Output file path (single mode) or output directory (batch mode).
                Defaults to the same location as the source PDF.
        dpi: Render resolution in dots per inch.
        jpeg_quality: JPEG compression quality (1–95).
        batch: If True, treat input as a directory and convert all .pdf files in it.
        with_notes: If True, extract and embed speaker notes from the matching
                    open Keynote document.
        keynote: Explicit path to the Keynote file to pull notes from.  When
                 provided, skips the interactive document picker.  Ignored in
                 batch mode (a picker runs per-file instead).
    """
    try:
        pdf_paths = validate_inputs(input, batch)
    except ValueError as exc:
        console.print(f"[bold red]Error:[/bold red] {exc}")
        raise SystemExit(1) from exc

    # Warn if --keynote is passed in batch mode (it doesn't apply there)
    if batch and keynote is not None:
        console.print(
            "[yellow]⚠ --keynote is ignored in batch mode. "
            "You will be prompted to pick a Keynote document for each PDF.[/yellow]"
        )
        keynote = None

    results: list[ConversionResult] = []

    # In single-file mode with notes, resolve the keynote path once up-front
    # so we don't interrupt the progress display mid-conversion.
    notes_mapping_single: NotesExtractionResult | None = None
    if with_notes and not batch:
        # Import locally to keep the top-level import surface small
        from keynote import pick_keynote_document  # type: ignore[import]

        resolved_keynote = keynote if keynote is not None else pick_keynote_document(pdf_paths[0])
        if resolved_keynote is None:
            console.print("[yellow]⚠ No Keynote document selected — proceeding without notes.[/yellow]")
        else:
            notes_mapping_single = _extract_notes_for_pdf(pdf_paths[0], resolved_keynote)

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        BarColumn(),
        TaskProgressColumn(),
        console=console,
        transient=False,
    ) as progress:
        overall_task = progress.add_task("[cyan]Converting PDFs…", total=len(pdf_paths))

        for pdf_path in pdf_paths:
            out_path = resolve_output_path(pdf_path, output, batch)

            # In batch mode with notes, resolve the keynote path per-file
            # (outside the progress block we already handled single-file mode).
            notes_mapping: NotesExtractionResult | None = notes_mapping_single
            if with_notes and batch:
                # Must stop progress briefly to allow interactive prompt
                progress.stop()
                try:
                    from keynote import pick_keynote_document  # type: ignore[import]

                    resolved_keynote = pick_keynote_document(pdf_path)
                    if resolved_keynote is None:
                        console.print(
                            f"[yellow]⚠ No Keynote document selected for {pdf_path.name} — skipping notes.[/yellow]"
                        )
                    else:
                        notes_mapping = _extract_notes_for_pdf(pdf_path, resolved_keynote)
                finally:
                    progress.start()

            # Open the PDF just long enough to get the page count for the progress bar.
            try:
                _doc = fitz.open(str(pdf_path))
                page_count = len(_doc)
                _doc.close()
            except Exception:
                page_count = None

            file_task = progress.add_task(
                f"  [dim]{pdf_path.name}[/dim]",
                total=page_count,
            )
            progress.update(overall_task, description=f"[cyan]{pdf_path.name}")

            result = convert_pdf_to_pptx(
                pdf_path=pdf_path,
                output_path=out_path,
                dpi=dpi,
                jpeg_quality=jpeg_quality,
                progress=progress,
                progress_task=file_task,
                notes_mapping=notes_mapping,
            )
            results.append(result)

            if result.status == ConversionStatus.FAILURE:
                console.print(
                    f"  [red]✗[/red] {pdf_path.name}: {result.error}",
                    highlight=False,
                )

            progress.update(file_task, visible=False)
            progress.advance(overall_task)

    print_summary(results)

    # Exit with a non-zero code if any conversion failed
    if any(r.status == ConversionStatus.FAILURE for r in results):
        raise SystemExit(1)


if __name__ == "__main__":
    auto_cli(main)
