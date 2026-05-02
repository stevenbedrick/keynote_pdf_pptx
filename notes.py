"""
notes.py — Extract speaker notes from a PPTX file and build a mapping from
Keynote slides to PDF pages.

The key challenge: Keynote can animate a single slide with multiple "click to
advance" build steps, which Keynote exports as separate PDF pages.  This
module reads the PPTX animation timing tree to count those steps per slide so
that it can compute the correct PDF-page range for each slide's notes.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

# Import models from the main module (they live there because ConversionResult
# also references them).
from pdf_to_pptx import NotesMismatchWarning, NotesExtractionResult, SlideNotesMapping

# PPTX presentation ML namespace
_PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ---------------------------------------------------------------------------
# Animation timing helpers
# ---------------------------------------------------------------------------


def count_click_build_steps(slide) -> int:  # noqa: ANN001
    """Count the number of click-advance build steps in *slide*'s timing tree.

    A "click-advance" step is a ``<p:par>`` element that is a direct child of
    the ``<p:childTnLst>`` of the ``<p:cTn nodeType="mainSeq">`` element, and
    whose ``<p:cTn><p:stCondLst>`` contains a ``<p:cond delay="indefinite"/>``.

    Returns 0 if no timing tree is found or if there are no such elements.
    """
    try:
        sp_tree = slide._element  # noqa: SLF001
    except AttributeError:
        return 0

    # Find the <p:cTn nodeType="mainSeq"> element
    main_seq_nodes = sp_tree.findall(
        f".//{{{_PPTX_NS}}}cTn[@nodeType='mainSeq']",
    )
    if not main_seq_nodes:
        return 0

    main_seq = main_seq_nodes[0]

    # The direct <p:childTnLst> of main_seq
    child_tn_lst = main_seq.find(f"{{{_PPTX_NS}}}childTnLst")
    if child_tn_lst is None:
        return 0

    click_steps = 0
    for par in child_tn_lst:
        if par.tag != f"{{{_PPTX_NS}}}par":
            continue
        # Look for <p:cTn><p:stCondLst><p:cond delay="indefinite"/>
        ctn = par.find(f"{{{_PPTX_NS}}}cTn")
        if ctn is None:
            continue
        st_cond_lst = ctn.find(f"{{{_PPTX_NS}}}stCondLst")
        if st_cond_lst is None:
            continue
        for cond in st_cond_lst:
            if cond.get("delay") == "indefinite":
                click_steps += 1
                break  # only need one match per <p:par>

    return click_steps


# ---------------------------------------------------------------------------
# Notes extraction
# ---------------------------------------------------------------------------


def extract_notes_paragraphs(slide) -> list[str]:  # noqa: ANN001
    """Return the speaker notes of *slide* as a list of paragraph strings.

    Each entry is the text of one paragraph (preserving empty strings for
    blank lines).  Returns ``[]`` if the slide has no notes or if an error
    occurs.
    """
    try:
        tf = slide.notes_slide.notes_text_frame
        return [para.text for para in tf.paragraphs]
    except Exception:  # noqa: BLE001
        return []


# ---------------------------------------------------------------------------
# Mapping builder
# ---------------------------------------------------------------------------


def build_notes_mapping(
    temp_pptx_path: Path,
    actual_pdf_page_count: int,
) -> NotesExtractionResult:
    """Build a :class:`NotesExtractionResult` from a temp PPTX export.

    Opens the PPTX at *temp_pptx_path*, iterates its slides (skipping hidden
    ones), and for each visible slide:

    - Counts the click-advance build steps to determine how many PDF pages it
      occupies (``steps + 1``).
    - Records the ``pdf_page_start`` / ``pdf_page_end`` range.
    - Extracts speaker-note paragraphs.

    Then compares the expected total page count to *actual_pdf_page_count* and
    sets the appropriate :class:`NotesMismatchWarning`.
    """
    prs = Presentation(str(temp_pptx_path))
    slides = prs.slides

    keynote_slide_count = len(slides)
    mappings: list[SlideNotesMapping] = []
    current_pdf_page = 0

    for slide_index, slide in enumerate(slides):
        # Skip hidden slides — they don't appear in the PDF export
        if slide._element.get("show") == "0":  # noqa: SLF001
            continue

        steps = count_click_build_steps(slide)
        pages_for_slide = steps + 1

        pdf_page_start = current_pdf_page
        pdf_page_end = current_pdf_page + pages_for_slide - 1

        notes = extract_notes_paragraphs(slide)

        mappings.append(
            SlideNotesMapping(
                keynote_slide_index=slide_index,
                pdf_page_start=pdf_page_start,
                pdf_page_end=pdf_page_end,
                notes_paragraphs=notes,
            )
        )

        current_pdf_page += pages_for_slide

    expected_pdf_pages = current_pdf_page  # == sum of all pages_for_slide

    if expected_pdf_pages == actual_pdf_page_count:
        warning = NotesMismatchWarning.NONE
        warning_message: str | None = None
    else:
        warning = NotesMismatchWarning.COUNT_MISMATCH
        warning_message = (
            f"Expected {expected_pdf_pages} PDF page(s) based on Keynote animation "
            f"build counts, but the PDF has {actual_pdf_page_count} page(s). "
            "Notes may be assigned to incorrect slides."
        )

    return NotesExtractionResult(
        mappings=mappings,
        keynote_slide_count=keynote_slide_count,
        expected_pdf_pages=expected_pdf_pages,
        actual_pdf_pages=actual_pdf_page_count,
        warning=warning,
        warning_message=warning_message,
    )
